from schemas import SlideContent, PosterExtraction
from pptx import Presentation

def map_to_slide_content(extracted: PosterExtraction, template: Presentation) -> SlideContent:
    placeholder_names = [shape.name.lower() for shape in template.slides[0].shapes if shape.has_text_frame]
    content_dict = {}
    for section in extracted.sections:
        key = section.heading.lower().replace(" ", "_")
        if key in placeholder_names:
            content_dict[key] = section.content
    return SlideContent(slide_type="summary_slide", title="Auto Generated Slide", content=content_dict)
