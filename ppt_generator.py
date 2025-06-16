import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ================= Helpers =================

def flatten_sections(sections):
    flattened = []
    for sec in sections:
        truncated_heading = truncate_and_wrap_heading(sec.heading)
        flattened.append((truncated_heading, sec.bullets))
    return flattened

def wrap_title(title, word_limit=9):
    words = title.split()
    if len(words) <= word_limit:
        return title
    first_line = " ".join(words[:word_limit])
    second_line = " ".join(words[word_limit:])
    return first_line + "\n" + second_line

def adaptive_title_font_size(title):
    word_count = len(title.split())
    if word_count <= 9:
        return 28
    elif word_count <= 15:
        return 24
    elif word_count <= 20:
        return 20
    else:
        return 18

def truncate_and_wrap_heading(heading, truncate_limit=15, wrap_limit=9):
    words = heading.split()
    words = words[:truncate_limit]
    if len(words) <= wrap_limit:
        return " ".join(words)
    first_line = " ".join(words[:wrap_limit])
    second_line = " ".join(words[wrap_limit:])
    return first_line + "\n" + second_line

def chunk_sections(flat_sections, chunk_size):
    all_pairs = []
    for heading, bullets in flat_sections:
        for bullet in bullets:
            all_pairs.append((heading, bullet))
    return [all_pairs[i:i + chunk_size] for i in range(0, len(all_pairs), chunk_size)]

# ================= Main PPT Generator =================

def generate_hierarchical_ppt(summary, template_path, output_path, image_path):
    prs = Presentation(template_path)
    flat_sections = flatten_sections(summary.sections)
    chunks = chunk_sections(flat_sections, chunk_size=5)

    for idx, chunk in enumerate(chunks):
        if idx == 0:
            slide = prs.slides[0]
        elif idx == 1:
            slide = prs.slides[1]
        else:
            base_slide = prs.slides[1]
            slide = prs.slides.add_slide(base_slide.slide_layout)
            for shape in base_slide.shapes:
                new_el = copy.deepcopy(shape.element)
                slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if 'Title Placeholder' in shape.name or 'Title Placeholder 2' in shape.name:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = wrap_title(summary.title)
                p.font.size = Pt(adaptive_title_font_size(summary.title))
                p.font.name = 'Calibri'
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)

            elif 'Content Placeholder' in shape.name or 'Content Placeholder 2' in shape.name:
                shape.text_frame.clear()
                last_heading = None
                for h, b in chunk:
                    if h != last_heading:
                        p = shape.text_frame.add_paragraph()
                        p.text = h
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.font.name = 'Calibri'
                        p.font.color.rgb = RGBColor(0, 0, 0)
                        last_heading = h
                    p = shape.text_frame.add_paragraph()
                    p.text = f"• {b}"
                    p.font.size = Pt(16)
                    p.font.bold = False
                    p.font.name = 'Calibri'
                    p.font.color.rgb = RGBColor(0, 0, 0)

        if idx == 0:
            for shape in slide.shapes:
                if "Image Placeholder 1" in shape.name or "Image Placeholder 2" in shape.name:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    slide.shapes.add_picture(image_path, left, top, width, height)

    prs.save(output_path)
    print(f"✅ Generated: {output_path}")
